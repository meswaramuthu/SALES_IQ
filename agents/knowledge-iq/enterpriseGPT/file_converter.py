"""Convert unsupported Office file types to plain text for Gemini.

Gemini/Vertex AI does not accept Office documents as inline data.
This module extracts plain text so the agent can process file attachments
uploaded via the ADK web UI.

Supported conversions:
  .docx / .doc   → python-docx
  .xlsx / .xls   → openpyxl
  .pptx / .ppt   → python-pptx
"""
from __future__ import annotations

import base64
import io
import logging
from typing import Optional

logger = logging.getLogger(__name__)

_MIME_TO_KIND: dict[str, str] = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/msword": "docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.ms-excel": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint": "pptx",
}


def _to_bytes(data) -> bytes:
    if isinstance(data, bytes):
        return data
    if isinstance(data, str):
        return base64.b64decode(data)
    return bytes(data)


def _extract_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            paragraphs.append("\t".join(c.text for c in row.cells))
    return "\n".join(paragraphs)


def _extract_xlsx(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"=== Sheet: {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            row_str = "\t".join("" if c is None else str(c) for c in row)
            if row_str.strip():
                lines.append(row_str)
    return "\n".join(lines)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"=== Slide {i} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


_EXTRACTORS = {
    "docx": _extract_docx,
    "xlsx": _extract_xlsx,
    "pptx": _extract_pptx,
}


def convert_to_text(mime_type: str, data) -> Optional[str]:
    """Extract plain text from an Office file.

    Returns the extracted text, or None if the MIME type is natively supported
    by Gemini and should be passed through unchanged.
    """
    kind = _MIME_TO_KIND.get(mime_type)
    if not kind:
        return None
    try:
        text = _EXTRACTORS[kind](_to_bytes(data))
        logger.info("Converted %s (%s) to %d chars of text", mime_type, kind, len(text))
        return text or "[File is empty or has no extractable text]"
    except Exception as exc:
        logger.warning("File conversion failed for %s: %s", mime_type, exc)
        return f"[Could not extract text from {mime_type}: {exc}]"
